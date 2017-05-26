#!/usr/bin/env python

'''
Convert lecture/slides from Slidoc Markdown to pptx format.

'''

from __future__ import print_function

import base64
import io
import json
import os
import re
import sys
import zipfile

try:
    import pptx
except ImportError:
    raise Exception('To read/write Powerpoint files, please install python package "pptx" using "pip install python-pptx"')

from pptx.util import Inches, Pt

import md2md


class MDParser(object):
    newline_norm_re =  re.compile(r'\r\n|\r')
    defaults_re =      re.compile(r'^ {0,3}(<!--slidoc-(defaults|options)\s+(.*?)-->|Slidoc:\s+(.*?))\s*$')
    header_re =        re.compile(r'^##[^#]')
    split_re =         re.compile(r'^(Extra|Notes|PluginDef|Tags):')
    notes_re =         re.compile(r'^Notes: *(.*)$')
    data_url_re =      re.compile(r'^data:([^;]+/[^;]+);base64,(.*)$')
    ref_def_re =       re.compile(r'''(^|\n) {0,3}\[([^\]]+)\]: +(\S+)( *\(.*\)| *'.*'| *".*")? *(?=\n|$)''')
    reflink_re =       re.compile(
        r'^ {0,3}!?\[('
        r'(?:\[[^^\]]*\]|[^\[\]]|\](?=[^\[]*\]))*'
        r')\]\s*\[([^^\]]*)\] *$'
    )

    rules_re = [ ('external_link', re.compile( r'''^ {0,3}(!?)\[([^\]]+)\]\(\s*(<)?(.*?)(?(3)>)(?:\s+['"](.*?)['"])?\s*\) *(\n|$)''') ),
                 ('hrule',         re.compile( r'^([-]{3,}) *(?:\n+|$)') )
                ]

    def __init__(self, args_dict={}):
        self.args_dict = args_dict
        self.images_zipfile = None
        self.images_map = {}
        self.slide_buffer = []
        self.filedir = ''
        self.new_slide()

    def parse_md(self, content, filehandle, filepath='', images_zip_handle=None):
        if images_zip_handle:
            self.images_zipfile = zipfile.ZipFile(images_zip_handle)
            self.images_map = dict( (os.path.basename(fpath), fpath) for fpath in self.images_zipfile.namelist() if os.path.basename(fpath))

        self.slide_start = True
        if filepath:
            self.filedir = os.path.dirname(os.path.realpath(filepath))

        content = self.newline_norm_re.sub('\n', content) # Normalize newlines

        self.ref_defs = {}
        for rmatch in self.ref_def_re.finditer(content):
            key = rmatch.group(2)
            link = rmatch.group(3)
            title = rmatch.group(4)[2:-1] if rmatch.group(4) else ''
            umatch = self.data_url_re.match(link)
            if not umatch:
                raise Exception('Invalid data URL for link '+link)
            content_type = umatch.group(1)
            self.ref_defs[key] = (title, content_type, base64.b64decode(umatch.group(2)) )

        while content:
            matched = None
            for rule_name, rule_re in self.rules_re:
                # Find the first match
                matched = rule_re.match(content)
                if matched:
                    break

            if matched:
                # Strip out matched text
                content = content[len(matched.group(0)):]

                if rule_name == 'external_link':
                    self.external_link(matched.group(0), matched.group(2), matched.group(4), matched.group(5) or '')

                elif rule_name == 'hrule':
                    self.hrule(matched.group(1))
                else:
                    raise Exception('Unknown rule: '+rule_name)

            else:
                # Strip out single line
                line, sep, content = content.partition('\n')
                if not line:
                    # Empty line
                    self.append_line(line, sep=sep)
                    continue

                if not self.slide_start and self.header_re.match(line):
                    self.new_slide()
                
                self.slide_start = False

                dmatch = self.defaults_re.match(line)
                if dmatch:
                    # Defaults
                    if len(self.slide_buffer) == 1:
                        # First slide
                        self.cur_slide['defaults'] = 'Slidoc: ' + (dmatch.group(3) or dmatch.group(4) or '')
                    continue
                    
                imatch = self.reflink_re.match(line)
                if imatch:
                    # Internal ref link
                    self.internal_link(imatch)
                    continue

                rmatch = self.ref_def_re.match(line)
                if rmatch:
                    # Ref definition
                    continue

                nmatch = self.notes_re.match(line)
                if nmatch:
                    # Notes section
                    if self.cur_slide['notes'] is None:
                        self.cur_slide['notes'] = []
                        line = nmatch.group(1)

                self.append_line(line, sep=sep)

        self.prs = pptx.Presentation()

        for slide in self.slide_buffer:
            self.dump_slide(slide)

        self.prs.save(filehandle)

    def new_slide(self):
        self.cur_slide = {'image': '', 'text': [], 'notes': None, 'defaults': ''}
        self.slide_buffer += [ self.cur_slide ]

    def append_line(self, text, sep='\n'):
        # Append line
        if self.cur_slide['notes'] is not None:
            self.cur_slide['notes'].append(text+sep)
        else:
            self.cur_slide['text'].append(text+sep)
        
    def hrule(self, text):
        self.new_slide()
        self.slide_start = True

    def dump_slide(self, slide):
        lines = slide['text']
        notes = slide['notes']

        blank_slide_layout = self.prs.slide_layouts[6]
        ppt_slide = self.prs.slides.add_slide(blank_slide_layout)

        ppt_notes = []
        if slide['defaults']:
            # First slide with defaults
            ppt_notes = [slide['defaults']+'\n']

        if slide['image']:
            # Image slide
            img_path = slide['image']
            img_name = os.path.basename(img_path)
            if img_path.startswith('#'):
                key = img_path[1:]
                if key not in self.ref_defs:
                    raise Exception('Invalid internal image reference: '+key)
                title, content_type, data = self.ref_defs[key]
                img_handle = io.BytesIO(data)
            elif self.images_zipfile:
                if img_name in self.images_map:
                    img_handle = io.BytesIO(self.images_zipfile.read(self.images_map[img_name]))
                else:
                    raise Exception('Image %s not found in zip archive' % img_name)
            else:
                img_handle = img_path

            top = Inches(1)
            left = Inches(0.5)
            height = Inches(5.5)
            pic = ppt_slide.shapes.add_picture(img_handle, left, top, height=height)
            ppt_notes += ['Slidoc:\n'] + lines
            lines = []

        else:
            # Text box slide
            shapes = ppt_slide.shapes

            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5.5)
            txBox = shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame

            offset = len(lines)
            for j, line in enumerate(lines):
                if line.startswith('<script>') or self.split_re.match(line):
                    # Split scripts/Extra/Notes/PluginDef/Tags off into notes portion (for readability)
                    offset = j
                    break
            
            self.add_text_to_frame(''.join(lines[:offset]), tf)
            ppt_notes += lines[offset:]

        if notes is not None:
            if ppt_notes:
                ppt_notes += ['\n']
            if notes[0].strip():
                ppt_notes += ['Notes: '+notes[0]]
            else:
                ppt_notes += ['Notes:\n']
            ppt_notes += notes[1:]

        if ppt_notes:
            notes_slide = ppt_slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            self.add_text_to_frame(''.join(ppt_notes), notes_text_frame)

    def add_text_to_frame(self, text, text_frame):
        # Normalize paragraph breaks and insert text into frame
        text_frame.word_wrap = True
        text = re.sub(r'\n +(?=\n)', r'\n', text)
        text = re.sub(r'\n{3,}', r'\n\n', text)
        paras = text.strip().split('\n\n')
        para_frame = text_frame
        while paras:
            para_frame.text = paras.pop(0)+'\n'
            if paras:
                para_frame = text_frame.add_paragraph()

    def internal_link(self, match):
        line = match.group(0)
        if line.lstrip().startswith("!"):
            text = match.group(1)
            key = match.group(2)
            if key not in self.ref_defs:
                raise Exception('Invalid internal image reference: '+key)
            title, content_type, data = self.ref_defs[key]
            if self.cur_slide['image']:
                raise Exception('Cannot handle more than one image per slide: '+line)
            self.cur_slide['image'] = '#'+key
            if title:
                self.append_line('![image1](%s "%s")' % (key, title) )
            else:
                self.append_line('![image1](%s)' % key )
        else:
            self.append_line(line)

    def external_link(self, line, text, link, title):
        if line.lstrip().startswith('!'):
            fpath = link
            if self.filedir:
                fpath = self.filedir + '/' + fpath
            bname, extn = os.path.splitext(os.path.basename(fpath))
            extn = extn.lower()
            if extn in ('.gif', '.jpg', '.jpeg', '.png', '.svg'):
                if self.cur_slide['image']:
                    raise Exception('Cannot handle more than one image per slide: '+line)
                self.cur_slide['image'] = fpath
                if title:
                    self.append_line('![image1](%s "%s")' % (bname, title))
                else:
                    self.append_line('![image1](%s)' % bname)
            else:
                self.append_line(line)
        else:
            self.append_line(line)


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Convert from Markdown to Powerpoint format')
    parser.add_argument('--overwrite', help='Overwrite files', action="store_true")
    parser.add_argument('file', help='Markdown filename', type=argparse.FileType('r'), nargs=argparse.ONE_OR_MORE)
    cmd_args = parser.parse_args()

    md_parser = MDParser(vars(cmd_args))

    fnames = []
    for f in cmd_args.file:
        fcomp = os.path.splitext(os.path.basename(f.name))
        fnames.append(fcomp[0])
        if fcomp[1] != '.md':
            sys.exit('Invalid file extension for '+f.name)

        if os.path.exists(fcomp[0]+'.pptx') and not cmd_args.overwrite:
            sys.exit("File %s.pptx already exists. Delete it or specify --overwrite" % fcomp[0])

    for j, f in enumerate(cmd_args.file):
        fname = fnames[j]
        outname = fname+".pptx"
        md_text = f.read()
        f.close()
        md_parser.parse_md(md_text, outname, f.name)

        print("Created ", outname, file=sys.stderr)
            
