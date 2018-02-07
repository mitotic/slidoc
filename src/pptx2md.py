#!/usr/bin/env python

'''
Convert lecture/slides from pptx to Slidoc Markdown format.

Use blank lines between list elements

If --embed_slides, embed whole slide as an image by default,
unless the string 'Answer:' occurs in the text of the slide at the start of a line.

All the text in the notes portion of a Powerpoint slide is appended to the text in the main slide.
(Explicitly use Notes: to delineate the notes portion of a Slidoc slide.)

If a line beginning with 'Slidoc:' is encountered and there is text below that line, all the text prior to that line is
essentially ignored and saved as the 'Extra:' portion of the slide. All the following text is retained.
For the first slide, any options after Slidoc: are retained. This allows multiple versions of a slide: a ppt version and a slidoc version
 
Use ![image0](image_name) to embed the slide itself as an image (after saving slides as .png/.jpg/.jpeg)
Use ![image1](image_name 'height=200') to embed specific image from main slide (can also specify height in title)
Note: image_name is optional in the above and may be omitted if using a default name derived from the file name.
      If specified, image_name should be unique and should not include the extension.
Use SlideContinue: as the last line of text to continue current Slidoc slide into next pptx slide (to embed multiple slide images)

If --auto_titles, the first line of text in the slide is used as the title, if it is the only line or if its followed by a blank line. (To have a slide without a title, split the text into at least two lines.)

'''

from __future__ import print_function

import base64
import io
import json
import os
import re
import sys
import zipfile

try:
    import pptx
    import PIL
except ImportError:
    raise Exception('To read/write Powerpoint files, please install python package "pptx" using "pip install python-pptx"')

import md2md


class PPTXParser(object):
    defaults_re = re.compile(r'^ {0,3}(<!--slidoc-(defaults|options)\s+(.*?)-->|Slidoc:\s+(.*?))\s*$')
    choice_re   = re.compile(r'^([A-P]\.\.)', re.IGNORECASE)
    image_re    = re.compile(r'''^ {0,3}!\[image(\d+)\]\(\s*([^'"]*?)(?:\s*(['"].*?['"]))?\s*\) *$''')
    notes_re    = re.compile(r'^Notes: *(.*)$')
    break_start_re = re.compile(r'''^[\w`'"]''')
    break_end_re   = re.compile(r'''^.*[\w`'"]$''')
    default_img_height = 440

    rules_re = [ ('external_link', re.compile( r'''^ {0,3}(!?)\[([^\]]+)\]\(\s*(<)?([\s\S]*?)(?(3)>)(?:\s+['"]([\s\S]*?)['"])?\s*\) *(\n|$)''') )
                ]

    def __init__(self, args_dict={}):
        self.args_dict = args_dict
        self.img_height = self.args_dict.get('img_height', self.default_img_height)
        self.nofile = False
        self.slide_zip_file = None
        self.prs = None
        self.image_keys = set()
        self.image_defs = []

    def parse_pptx(self, filehandle, filepath='', slides_zip_handle=None, nofile=False):
        # Returns (md_text, images_zip or None)
        # nofile if not actualy reading a file from disk
        self.nofile = nofile
        if slides_zip_handle:
            self.slide_zip_file = zipfile.ZipFile(slides_zip_handle)

        self.prs = pptx.Presentation(filehandle)

        self.filename = os.path.splitext(os.path.basename(filepath))[0] if filepath else ''
        if filepath and not self.nofile:
            self.filedir = os.path.dirname(os.path.realpath(filepath))
            self.fileprefix = self.filename
            self.file_mtime = os.path.getmtime(filepath) if os.path.exists(filepath) else 0
        else:
            self.filedir = ''
            self.fileprefix = ''
            self.file_mtime = 0
        self.img_dir = self.args_dict.get('img_dir','')
        self.img_bytes = None
        self.img_zip = None
        if self.img_dir is None:
            self.img_dir = self.fileprefix+'_images' if self.fileprefix else '_images'
        elif self.img_dir.endswith('.zip'):
            self.img_bytes = io.BytesIO()
            self.img_zip = zipfile.ZipFile(self.img_bytes, 'w')

        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        slide_buffer = []
        prev_continued = False

        nslides = len(self.prs.slides)
        self.img_count = 0
        for slide_indx, slide in enumerate(self.prs.slides):
            slide_num = slide_indx+1
            full_slide = ('%s/Slide%02d' if nslides >= 10 else '%s/Slide%d') % (self.filename, slide_num)
            slide_images = []
            slide_text = ''
            shape_list = [(shape.top, shape) for shape in slide.shapes]
            shape_list.sort()
            for _, shape in shape_list:
                if isinstance(shape, pptx.shapes.picture.Picture):
                    # Image
                    image_obj = shape.image
                    image_blob = image_obj.blob
                    img_width, img_height = image_obj.size

                    if shape.crop_left or shape.crop_right or shape.crop_top or shape.crop_bottom:
                        # Crop image (pixel coordinate system origin top-left: refers to pixel corners, not centers)
                        c_left   = int(round(img_width*shape.crop_left))
                        c_right  = int(round(img_width*(1.0-shape.crop_right)))
                        c_top    = int(round(img_height*shape.crop_top))
                        c_bottom = int(round(img_height*(1.0-shape.crop_bottom)))
                        in_stream = io.BytesIO(image_blob)
                        try:
                            pil_image = PIL.Image.open(in_stream)
                            format = pil_image.format
                            pil_image = pil_image.crop( (c_left, c_top, c_right, c_bottom) )
                            in_stream.close()
                            out_stream = io.BytesIO()
                            pil_image.save(out_stream, format=format)
                            image_blob = out_stream.getvalue()
                            img_width, img_height = pil_image.size
                            out_stream.close()
                        except Exception, excp:
                            print('pptx2md: Error: in processing shape image %s: %s' % (image_obj.filename, excp))
                            continue

                    slide_images.append( {'ext': image_obj.ext, 'width': img_width, 'height': img_height, 'blob': image_blob} )
                    
                    if self.args_dict.get('img_inline'):
                        slide_text += "\n![image%d]()\n\n" % len(slide_images)
                elif shape.has_text_frame:
                    # Text
                    if slide_text:
                        slide_text += '\n'
                    prev_level = 0
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ''
                        fenced = False
                        for run in paragraph.runs:
                            run_text = run.text
                            if run_text.startswith('```'):
                                fenced = not fenced
                            if para_text and run_text:
                                if fenced or run_text.startswith('```') or self.args_dict.get('line_breaks'):
                                    para_text += '\n'
                                elif self.break_end_re.match(para_text) and self.break_start_re.match(run_text):
                                    para_text += ' '
                            para_text += run_text

                        para_text = md2md.restore_angular(md2md.asciify(para_text))

                        if self.choice_re.match(para_text):
                            # Choice option
                            slide_text = self.force_para_break(slide_text)
                        if paragraph.level and para_text:
                            slide_text += ''.join([' ']*(paragraph.level-1)) + '- '
                        elif prev_level:
                            slide_text += '\n'
                        prev_level = paragraph.level
                        slide_text += para_text + '\n\n'

            if slide.has_notes_slide:
                # Append all notes text to slide text (unless --notes is set)
                notes_text = slide.notes_slide.notes_text_frame.text
                if self.args_dict.get('notes'):
                    if 'Notes:' not in slide_text and 'Notes:' not in notes_text and not notes_text.startswith('Slidoc:'):
                        # Treat all notes as notes
                        notes_text = 'Notes:\n' + notes_text
                slide_text = self.force_para_break(slide_text) + notes_text

            # Slidoc
            defaults_text = ''
            extra_text = ''
            continued = False
            embed_slide = False

            if slide_text.rstrip().endswith('SlideContinue:'):
                slide_text = slide_text.rstrip()[:-len('SlideContinue:')].rstrip()
                continued = True

            elif self.args_dict.get('embed_slides') and '\nAnswer:' not in slide_text:
                # Insert entire slide as image by default
                embed_slide = True

            elif not slide_text.strip():
                # No text 
                if not slide_images:
                    # No images
                    if not continued and slide_num != nslides:
                        slide_buffer.append('\n---\n\n')
                    continue

            has_notes = '\nNotes:' in slide_text
            slide_lines = slide_text.strip().split('\n')

            if self.args_dict.get('auto_titles') and not prev_continued and (len(slide_lines) == 1 or (len(slide_lines) > 1 and not slide_lines[1].strip())) and slide_lines[0][0] not in ' #!' and slide_lines[0].strip()[-1] not in '.?:':
                # Automatic titles (use space in first line to suppress it)
                title = slide_lines[0].strip()
                slide_lines = slide_lines[1:]
                if slide_lines and not slide_lines[0].strip():
                    slide_lines = slide_lines[1:]
                if slide_num == 1:
                    title = '# ' + title
                else:
                    title = '## ' + title
                slide_lines = [title, ''] + slide_lines

            if embed_slide:
                # Embedding slide; save slide text as Extra:
                if slide_lines[0].startswith('##') and not has_notes:
                    # Save slide title in Notes:
                    extra_text = '\n'.join(slide_lines[1:])
                    slide_lines = ['', 'Notes:'] + slide_lines[0:1]
                else:
                    extra_text = '\n'.join(slide_lines)
                    slide_lines = []
                # Insert entire slide as image
                slide_lines = [ '![image0]()' ] + slide_lines

            images_copied = set()
            md_lines = []
            for line_num, line in enumerate(slide_lines):
                dmatch = self.defaults_re.match(line)
                if dmatch:
                    # Settings line
                    if slide_num == 1 and not defaults_text:
                        # First slide; retain default settings
                        defaults_text = dmatch.group(3) or dmatch.group(4) or ''
                    if any(x.strip() for x in slide_lines[line_num+1:]):
                        # "discard" all text prior to Slidoc: (i.e., save as extra) and use text following Slidoc: for slide
                        extra_text = ''.join(md_lines)
                        md_lines = []
                    continue

                lmatch = self.image_re.match(line)
                if lmatch:
                    # Image
                    image_num = int(lmatch.group(1))
                    img_name = lmatch.group(2).strip()
                    img_params = lmatch.group(3) or ''
                    images_copied.add(image_num)
                    if image_num == 0:
                        # Slide image: ![image0]()
                        iheight = self.img_height
                        if not img_params:
                            img_params = "'height=%d'" % iheight
                        img_path = self.slide_image_path(full_slide)
                        img_ref = self.copy_image(image_num, self.read_file(img_path), img_params=img_params, img_name=img_name, img_path=img_path)
                    elif image_num <= len(slide_images):
                        # Embedded image: ![image1]()
                        slide_image = slide_images[image_num-1]
                        fwidth = self.compute_img_fracwidth(slide_image['width'], slide_image['height'], len(slide_images))
                        if not img_params:
                            img_params = "'width=%s'" % fwidth
                        img_ref = self.copy_image(image_num, slide_image['blob'], img_params=img_params, img_name=img_name, img_ext=slide_image['ext'])
                    else:
                        img_ref = 'Missing/annotated image%d(%s)' % (image_num, img_name or '')
                        print('pptx2md: Error: ![image%d]() in slide %d exceeds available image count of %d. Only raw images are recognized in slides. Remove any annotation and re-paste raw image.' % (image_num, slide_num, len(slide_images)), file=sys.stderr)
                    md_lines.append('\n'+img_ref+'\n\n')
                else:
                    # Not image
                    md_lines.append(line+'\n')

            if not images_copied and slide_images:
                # Copy embedded images by default
                md_images = ["\n"]
                for image_num, slide_image in enumerate(slide_images):
                    fwidth = self.compute_img_fracwidth(slide_image['width'], slide_image['height'], len(slide_images))
                    img_params = "'width=%s'" % fwidth
                    img_ref = self.copy_image(image_num+1, slide_image['blob'], img_params=img_params, img_ext=slide_image['ext'])
                    md_images.append('\n'+img_ref+'\n\n')
                offset = len(md_lines)
                for line_num, md_line in enumerate(md_lines):
                    if md_line.startswith('Answer:'):
                        # Insert images before Answer:
                        offset = line_num
                        break
                md_lines = md_lines[:offset] + md_images + md_lines[offset:]
 
            md_text = self.force_para_break(''.join(md_lines))
            if defaults_text:
                md_text = 'Slidoc: '+defaults_text + '\n' + md_text
            if extra_text:
                md_text += 'Extra:\n'+extra_text+'\n\n'

            if not continued and slide_num != nslides:
                # End of slide
                md_text += '---\n\n'

            slide_buffer.append(md_text)
            prev_continued = continued

        all_text = ''.join(slide_buffer)
        if self.image_defs:
            all_text += '\n\n' + '\n\n'.join(self.image_defs) + '\n\n'

        zipped_md = None
        if self.img_zip:
            if self.args_dict.get('zip_md'):
                self.img_zip.writestr('content.md', all_text)
            self.img_zip.close()
            zipped_md = self.img_bytes.getvalue()
        return all_text, zipped_md

    def compute_img_fracwidth(self, width, height, nimages):
        tem_width = (3.0/4.0)*width/(1.0*height)
        perc = (100.0 if self.args_dict.get('expand_images') else 70.0) / nimages
        return '%d%%' % int( round(perc*min(1,tem_width)))

    def compute_img_height(self, width, height, max_height):
        max_width = int(max_height * (4.0/3.0))
        tem_height = int(height*max_width/(1.0*width))
        if width/(1.0*height) > 4.0/3.0:
            # Elongated
            return tem_height if self.args_dict.get('expand_images') else min(tem_height, max_height)
        else:
            return max_height if self.args_dict.get('expand_images') else min(height, max_height)

    def force_para_break(self, text):
        if not text.endswith('\n'):
            return text + '\n\n'
        elif not text[:-1].endswith('\n'):
            return text + '\n'
        else:
            return text

    def external_link(self, line, text, link, title):
        if line.lstrip().startswith('!'):
            fpath = link
            if self.filedir:
                fpath = self.filedir + '/' + fpath
            _, extn = os.path.splitext(os.path.basename(fpath))
            extn = extn.lower()
            if extn in ('.gif', '.jpg', '.jpeg', '.png', '.svg'):
                self.cur_slide['image'] = fpath
            else:
                self.cur_slide['text'].append(line)
        else:
            self.cur_slide['text'].append(line)

    def slide_image_path(self, image_file):
        extensions = ('.png', '.jpg', '.jpeg')
        for extn in extensions:
            ipath = ''
            if self.slide_zip_file:
                if image_file+extn in self.slide_zip_file.namelist():
                    ipath = image_file+extn
                    break
            elif not self.nofile and os.path.exists(image_file+extn):
                ipath = image_file+extn
                break
        if not ipath:
            raise Exception('pptx2md: Slide image file '+image_file+'/'.join(extensions)+' not found! Please export slideshow as images')

        if not self.slide_zip_file and os.path.getmtime(ipath) < self.file_mtime:
            raise Exception('pptx2md: Slide image file '+ipath+' older than slide file! Please re-export slideshow as images')

        return ipath

    def read_file(self, filepath):
        if self.slide_zip_file:
            return self.slide_zip_file.read(filepath)
        with open(filepath) as f:
            return f.read()

    def copy_image(self, image_num, img_data, img_params='', img_path='', img_ext='', img_name=''):
        # Return image reference
        if not img_ext and img_path:
            img_ext = os.path.splitext(os.path.basename(img_path))[1][1:].lower()

        if not self.args_dict.get('img_dir'):
            # Embed images within Markdown
            if img_name:
                key = img_name
            else:
                key = 'img%02d' % (len(self.image_defs)+1)
            if key in self.image_keys:
                raise Exception('Duplicate image file name: '+key)
            self.image_keys.add(key)
            ctype = 'jpeg' if img_ext == 'jpg' else img_ext
            self.image_defs.append('[%s]: data:image/%s;base64,%s %s' % (key, ctype, base64.b64encode(img_data), img_params) )
            return '![image%d][%s]' % (image_num, key)

        if img_name:
            img_copy = img_name + '.' + img_ext
        else:
            prefix = self.fileprefix + '-' if self.fileprefix else ''
            if img_path:
                img_copy = prefix + os.path.basename(img_path)
            else:
                self.img_count += 1
                img_copy = prefix + ('image%02d.%s' % (self.img_count, img_ext))

        if self.img_zip:
            # Write image file to zip archive
            zprefix = os.path.splitext(os.path.basename(self.img_dir))[0]
            if zprefix:
                img_copy = zprefix + '/' + img_copy
            self.img_zip.writestr(img_copy, img_data)
        elif not self.nofile:
            # Write image file to disk
            if self.img_dir:
                img_copy = self.img_dir + '/' + img_copy
                if not os.path.exists(self.img_dir):
                    os.makedirs(self.img_dir)

            with open(img_copy, 'w') as f:
                f.write(img_data)

        return '![image%d](%s %s)' % (image_num, img_copy, img_params)


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Convert from Powerpoint (pptx) to Markdown format')
    parser.add_argument('-a', '--auto_titles', help='Automatically generate titles', action="store_true")
    parser.add_argument('-e', '--embed_slides', help='Embed image of whole slide by default (unless Answer: is present)', action="store_true")
    parser.add_argument('-x', '--expand_images', help='Expand images to fill slide', action="store_true")
    parser.add_argument('-i', '--img_inline', help='Insert images inline', action="store_true")
    parser.add_argument('--img_dir', help='Image directory: omit to embed images; name.zip to create zip archive', default=None)
    parser.add_argument('--img_height', help='Image height', default=440)
    parser.add_argument('-l', '--line_breaks', help='Break lines within text box paragraphs', action="store_true")
    parser.add_argument('-n', '--notes', help='Treat slide notes as Notes: (instead of as text)', action="store_true")
    parser.add_argument('--overwrite', help='Overwrite files', action="store_true")
    parser.add_argument('--zip_md', help='Include md file in zip archive', action="store_true")
    parser.add_argument('file', help='Powerpoint filename', type=argparse.FileType('r'), nargs=argparse.ONE_OR_MORE)
    cmd_args = parser.parse_args()

    ppt_parser = PPTXParser(vars(cmd_args))

    fnames = []
    for f in cmd_args.file:
        fcomp = os.path.splitext(os.path.basename(f.name))
        fnames.append(fcomp[0])
        if fcomp[1] != '.pptx':
            sys.exit('Invalid file extension for '+f.name)

        if os.path.exists(fcomp[0]+'.md') and not cmd_args.overwrite:
            sys.exit("File %s.md already exists. Delete it or specify --overwrite" % fcomp[0])

    for j, f in enumerate(cmd_args.file):
        fname = fnames[j]
        outname = fname+".md"
        md_text, images_zip = ppt_parser.parse_pptx(f, f.name)
        with io.open(outname, 'w', encoding='utf8') as f:
            f.write(md_text)
        print("Created ", outname, file=sys.stderr)

        if images_zip:
            with io.open(cmd_args.img_dir, 'wb') as f:
                f.write(images_zip)
            print("Created ", cmd_args.img_dir, "with images", file=sys.stderr)

            
